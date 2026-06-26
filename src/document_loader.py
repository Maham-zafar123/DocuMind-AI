from __future__ import annotations

import io
from dataclasses import dataclass
from typing import List

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


@dataclass
class DocumentChunk:
    filename: str
    chunk_id: int
    text: str


def _read_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
        for page_no, page in enumerate(pdf, start=1):
            page_text = page.get_text("text")
            if page_text.strip():
                text_parts.append(f"\n[Page {page_no}]\n{page_text}")
    return "\n".join(text_parts)


def _read_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                tables.append(" | ".join(cells))
    return "\n".join(paragraphs + tables)


def _read_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides.append(f"\n[Slide {idx}]\n" + "\n".join(slide_text))
    return "\n".join(slides)


def _read_txt(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_text(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    file_bytes = uploaded_file.getvalue()
    if name.endswith(".pdf"):
        return _read_pdf(file_bytes)
    if name.endswith(".docx"):
        return _read_docx(file_bytes)
    if name.endswith(".pptx"):
        return _read_pptx(file_bytes)
    if name.endswith(".txt"):
        return _read_txt(file_bytes)
    raise ValueError(f"Unsupported file type: {uploaded_file.name}")


def chunk_text(filename: str, text: str, chunk_size: int, overlap: int) -> List[DocumentChunk]:
    clean = " ".join(text.split())
    if not clean:
        return []
    chunks = []
    start = 0
    chunk_id = 1
    while start < len(clean):
        end = min(start + chunk_size, len(clean))
        chunk = clean[start:end]
        chunks.append(DocumentChunk(filename=filename, chunk_id=chunk_id, text=chunk))
        chunk_id += 1
        if end == len(clean):
            break
        start = max(0, end - overlap)
    return chunks
